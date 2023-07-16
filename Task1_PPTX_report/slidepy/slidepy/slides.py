import io

import typeguard
from typing import TypedDict

import pptx.slide
from pptx.util import Inches

import numpy as np
import matplotlib.pyplot as plt

class Slide (pptx.slide.Slide):
    """
    Base class representing a slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        layout (pptx.slide.SlideLayout): The slide layout to use for the slide.
    """

    def __init__(self, presentation, layout):
        self.ref = presentation.slides.add_slide(layout)

    @classmethod
    def from_dict(cls, presentation, dict):
        """
        Create a slide object from a dictionary representation.

        Args:
            presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
            dict (dict): Dictionary representation of the slide.

        Returns:
            Slide: The created slide object.
        
        Raises:
            ValueError: If an invalid slide type is provided.
        """
        type = dict.pop("type")

        match type:
            case "title":
                return TitleSlide(presentation, **dict)
            case "text":
                return TextSlide(presentation, **dict)
            case "list":
                return ListSlide(presentation, **dict)
            case "picture":
                return PictureSlide(presentation, **dict)
            case "plot":
                return PlotSlide(presentation, **dict)
            case _:
                raise ValueError(f'"{type}" is not a valid slide type!')


class TitleSlide (Slide):
    """
    Class representing a title slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str, optional): The title of the slide. Defaults to "".
        content (str, optional): The content of the slide. Defaults to "".
    """

    @typeguard.typechecked
    def __init__(self, presentation, title:str="", content:str=""):
        layout = presentation.slide_layouts[0]
        super(TitleSlide, self).__init__(presentation, layout)
        self.ref.shapes.title.text = title
        self.ref.placeholders[1].text = content


class TextSlide (Slide):
    """
    Class representing a text slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str, optional): The title of the slide. Defaults to "".
        content (str, optional): The content of the slide. Defaults to "".
    """

    @typeguard.typechecked
    def __init__(self, presentation, title:str="", content:str=""):
        layout = presentation.slide_layouts[5]
        super(TextSlide, self).__init__(presentation, layout)
        self.ref.shapes.title.text = title
        left = width = height = Inches(1)
        top = Inches(2)
        txBox = self.ref.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = content


class ListElements(TypedDict):
    """
    Typed dictionary representing the elements of a list slide.

    Attributes:
        level (int): The level of the list item.
        text (str): The text of the list item.
    """
    level:int
    text:str


class ListSlide (Slide):
    """
    Class representing a list slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str, optional): The title of the slide. Defaults to "".
        content (list[ListElements], optional): The list items of the slide. Defaults to [].
    """

    @typeguard.typechecked
    def __init__(self, presentation, title:str="", content:list[ListElements]=[]):
        layout = presentation.slide_layouts[1]
        super(ListSlide, self).__init__(presentation, layout)
        shapes = self.ref.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame

        for kwargs in content:
            p = tf.add_paragraph()
            p.text = kwargs["text"]
            p.level = kwargs["level"]


class PictureSlide (Slide):
    """
    Class representing a picture slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str, optional): The title of the slide. Defaults to "".
        content (str, optional): The file path of the picture. Defaults to "".
    """

    @typeguard.typechecked
    def __init__(self, presentation, title:str="", content:str=""):
        layout = presentation.slide_layouts[5]
        super(PictureSlide, self).__init__(presentation, layout)
        self.ref.shapes.title.text = title
        left = Inches(1)
        top = Inches(2)
        height = Inches(5)

        img = self.ref.shapes.add_picture(content, left, top, height=height)


class PlotSlide (Slide):
    """
    Class representing a plot slide in the PowerPoint presentation.

    Args:
        presentation (pptx.presentation.Presentation): The PowerPoint presentation object.
        title (str, optional): The title of the slide. Defaults to "".
        content (str, optional): The file path of the data file. Defaults to "".
        configuration (dict, optional): The configuration for the plot. Defaults to {}.
    """

    @typeguard.typechecked
    def __init__(self, presentation, title:str="", content:str="", configuration:dict={}):
        layout = presentation.slide_layouts[5]
        super(PlotSlide, self).__init__(presentation, layout)
        self.ref.shapes.title.text = title
        left = Inches(1)
        top = Inches(2)
        height = Inches(5)

        data = np.loadtxt(content, delimiter=';')
        plt.plot(data[:,0], data[:,1], linewidth=2.0)
        plt.xlabel(configuration["x-label"])
        plt.ylabel(configuration["y-label"])

        img_stream = io.BytesIO()
        plt.savefig(img_stream)
        img = self.ref.shapes.add_picture(img_stream, left, top, height=height)
